#!/usr/bin/env python3
from __future__ import annotations

import csv
import os
import random
from copy import deepcopy
from pathlib import Path
from typing import Any, Dict, List, Literal, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

try:
    import tomllib  # Python 3.11+
except ModuleNotFoundError:
    import tomli as tomllib  # pip install tomli

BG = RGBColor(240, 247, 255)
BLUE = RGBColor(37, 99, 235)
TEXT = RGBColor(15, 23, 42)
SUBTEXT = RGBColor(71, 85, 105)
CORAL = RGBColor(249, 112, 102)
WHITE = RGBColor(255, 255, 255)

VocabPair = Tuple[str, str]
PrimarySide = Literal["term", "definition"]
BackgroundMode = Literal["fixed", "cycle"]

IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg"}

DEFAULT_STYLE_CONFIG: Dict[str, Any] = {
    "labels": {
        "game_title": "Vocabulary Games",
        "source_prefix": "Source:",
        "set_prefix": "Set",
        "vocabulary_suffix": "Vocabulary",
    },
    "palettes": {
        "colors": {
            "default": {
                "bg": "#F0F7FF",
                "blue": "#2563EB",
                "text": "#0F172A",
                "subtext": "#475569",
                "coral": "#F97066",
                "white": "#FFFFFF",
            }
        },
        "fonts": {
            "title_pill": {
                "name": "Aptos",
                "size": 16,
                "bold": True,
            },
            "title_main": {
                "name": "Aptos Display",
                "size": 28,
                "bold": True,
            },
            "title_section": {
                "name": "Aptos Display",
                "size": 24,
                "bold": True,
            },
            "body": {
                "name": "Aptos",
                "size": 16,
                "bold": False,
            },
            "vocab_primary": {
                "name": "Aptos Display",
                "size": 24,
                "bold": True,
            },
            "vocab_secondary": {
                "name": "Aptos",
                "size": 20,
                "bold": False,
            },
        },
    },
    "styles": {
        "title_slide": {
            "color_palette": "default",
            "pill_font": "title_pill",
            "main_font": "title_main",
            "section_font": "title_section",
            "body_font": "body",
            "overlay_transparency": 0.22,
            "show_card": True,
            "card_transparency": 0.18,
        },
        "vocab_slide": {
            "color_palette": "default",
            "primary_font": "vocab_primary",
            "secondary_font": "vocab_secondary",
            "overlay_transparency": 0.22,
            "show_card": True,
            "card_transparency": 0.18,
        },
    },
}


def deep_merge_dicts(base: Dict[str, Any], override: Dict[str, Any]) -> Dict[str, Any]:
    merged = deepcopy(base)
    for key, value in override.items():
        if key in merged and isinstance(merged[key], dict) and isinstance(value, dict):
            merged[key] = deep_merge_dicts(merged[key], value)
        else:
            merged[key] = value
    return merged


def get_default_style_config_path() -> Path:
    xdg_config_home = os.environ.get("XDG_CONFIG_HOME", "").strip()
    if xdg_config_home:
        config_home = Path(xdg_config_home).expanduser()
    else:
        config_home = Path.home() / ".config"
    return config_home / "mainkata" / "style.toml"


def resolve_style_config_path(
    style_config_file: str | Path | None = None,
) -> Path | None:
    if style_config_file is not None:
        config_path = Path(style_config_file).expanduser().resolve()
        if not config_path.exists():
            raise FileNotFoundError(f"Style config file not found: {config_path}")
        if not config_path.is_file():
            raise ValueError(f"Style config path is not a file: {config_path}")
        return config_path

    default_path = get_default_style_config_path().resolve()
    if default_path.exists() and default_path.is_file():
        return default_path

    return None


def load_style_config(style_config_file: str | Path | None = None) -> Dict[str, Any]:
    config_path = resolve_style_config_path(style_config_file)
    if config_path is None:
        return deepcopy(DEFAULT_STYLE_CONFIG)

    with config_path.open("rb") as f:
        loaded = tomllib.load(f)

    return deep_merge_dicts(DEFAULT_STYLE_CONFIG, loaded)


def hex_to_rgb_color(value: str) -> RGBColor:
    text = value.strip().lstrip("#")
    if len(text) != 6:
        raise ValueError(f"Invalid hex color: {value!r}")
    try:
        return RGBColor.from_string(text.upper())
    except ValueError as exc:
        raise ValueError(f"Invalid hex color: {value!r}") from exc


def resolve_color_palette(
    style_config: Dict[str, Any], palette_name: str
) -> Dict[str, RGBColor]:
    palettes = style_config["palettes"]["colors"]
    if palette_name not in palettes:
        raise ValueError(f"Unknown color palette: {palette_name}")
    palette = palettes[palette_name]
    return {key: hex_to_rgb_color(value) for key, value in palette.items()}


def resolve_font_palette(
    style_config: Dict[str, Any], font_name: str
) -> Dict[str, Any]:
    palettes = style_config["palettes"]["fonts"]
    if font_name not in palettes:
        raise ValueError(f"Unknown font palette: {font_name}")
    font = palettes[font_name]
    return {
        "name": str(font["name"]),
        "size": int(font["size"]),
        "bold": bool(font["bold"]),
    }


def resolve_title_slide_style(style_config: Dict[str, Any]) -> Dict[str, Any]:
    style = style_config["styles"]["title_slide"]
    return {
        "colors": resolve_color_palette(style_config, style["color_palette"]),
        "pill_font": resolve_font_palette(style_config, style["pill_font"]),
        "main_font": resolve_font_palette(style_config, style["main_font"]),
        "section_font": resolve_font_palette(style_config, style["section_font"]),
        "body_font": resolve_font_palette(style_config, style["body_font"]),
        "overlay_transparency": float(style["overlay_transparency"]),
        "show_card": bool(style["show_card"]),
        "card_transparency": float(style["card_transparency"]),
    }


def resolve_vocab_slide_style(style_config: Dict[str, Any]) -> Dict[str, Any]:
    style = style_config["styles"]["vocab_slide"]
    return {
        "colors": resolve_color_palette(style_config, style["color_palette"]),
        "primary_font": resolve_font_palette(style_config, style["primary_font"]),
        "secondary_font": resolve_font_palette(style_config, style["secondary_font"]),
        "overlay_transparency": float(style["overlay_transparency"]),
        "show_card": bool(style["show_card"]),
        "card_transparency": float(style["card_transparency"]),
    }


def resolve_csv_path(csv_file: str | Path) -> Path:
    csv_path = Path(csv_file).expanduser().resolve()
    if not csv_path.exists():
        raise FileNotFoundError(f"CSV file not found: {csv_path}")
    return csv_path


def resolve_output_path(csv_path: Path, output: str | Path | None = None) -> Path:
    if output:
        output_path = Path(output).expanduser().resolve()
    else:
        output_path = csv_path.with_name(csv_path.stem + "_vocab_sets.pptx")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    return output_path


def resolve_background_dir(background_dir: str | Path) -> Path:
    bg_dir = Path(background_dir).expanduser().resolve()
    if not bg_dir.exists():
        raise FileNotFoundError(f"Background directory not found: {bg_dir}")
    if not bg_dir.is_dir():
        raise ValueError(f"Background path is not a directory: {bg_dir}")
    return bg_dir


def validate_generation_options(
    set_count: int,
    set_size: int,
    primary_side: PrimarySide,
) -> None:
    if set_count < 1:
        raise ValueError("--sets must be at least 1.")
    if set_size < 1:
        raise ValueError("--set-size must be at least 1.")
    if primary_side not in {"term", "definition"}:
        raise ValueError("--primary-side must be either 'term' or 'definition'.")


def validate_background_options(
    background_dir: str | Path | None,
    background_mode: BackgroundMode,
    background_image_number: int | None,
    background_cycle_start: int | None,
    background_cycle_end: int | None,
) -> None:
    if background_dir is None:
        return

    if background_mode not in {"fixed", "cycle"}:
        raise ValueError("--background-mode must be 'fixed' or 'cycle'.")

    if background_mode == "fixed":
        if background_image_number is None or background_image_number < 1:
            raise ValueError(
                "--background-image-number must be >= 1 when "
                "--background-mode=fixed."
            )
        if background_cycle_start is not None or background_cycle_end is not None:
            raise ValueError(
                "--background-cycle-start and --background-cycle-end cannot be "
                "used with --background-mode=fixed."
            )

    if background_mode == "cycle":
        if background_image_number is not None:
            raise ValueError(
                "--background-image-number cannot be used with "
                "--background-mode=cycle."
            )

        one_range_value = (background_cycle_start is None) != (
            background_cycle_end is None
        )
        if one_range_value:
            raise ValueError(
                "--background-cycle-start and --background-cycle-end must be "
                "provided together."
            )

        if (
            background_cycle_start is not None
            and background_cycle_end is not None
            and background_cycle_start < 1
        ):
            raise ValueError("--background-cycle-start must be at least 1.")

        if (
            background_cycle_start is not None
            and background_cycle_end is not None
            and background_cycle_end < 1
        ):
            raise ValueError("--background-cycle-end must be at least 1.")

        if (
            background_cycle_start is not None
            and background_cycle_end is not None
            and background_cycle_start > background_cycle_end
        ):
            raise ValueError(
                "--background-cycle-start cannot be greater than "
                "--background-cycle-end."
            )


def validate_visual_options(
    title_slide_overlay_transparency: float,
    vocab_slide_overlay_transparency: float,
    title_card_transparency: float,
    vocab_card_transparency: float,
) -> None:
    if not 0.0 <= title_slide_overlay_transparency <= 1.0:
        raise ValueError(
            "--title-slide-overlay-transparency must be between 0.0 and 1.0."
        )
    if not 0.0 <= vocab_slide_overlay_transparency <= 1.0:
        raise ValueError(
            "--vocab-slide-overlay-transparency must be between 0.0 and 1.0."
        )
    if not 0.0 <= title_card_transparency <= 1.0:
        raise ValueError("--title-card-transparency must be between 0.0 and 1.0.")
    if not 0.0 <= vocab_card_transparency <= 1.0:
        raise ValueError("--vocab-card-transparency must be between 0.0 and 1.0.")


def is_valid_image(path: Path) -> bool:
    try:
        with Image.open(path) as im:
            im.verify()
        return True
    except Exception:
        return False


def list_background_images(background_dir: Path) -> List[Path]:
    candidates = sorted(
        p
        for p in background_dir.iterdir()
        if p.is_file() and p.suffix.lower() in IMAGE_SUFFIXES
    )

    bad: List[Path] = []
    images: List[Path] = []

    for p in candidates:
        if is_valid_image(p):
            images.append(p)
        else:
            bad.append(p)

    if bad:
        bad_list = "\n".join(str(p) for p in bad)
        raise ValueError(
            "The following files in the background directory are not valid PNG/JPG "
            "images or use unsupported encodings:\n\n"
            f"{bad_list}\n\n"
            "Please remove or convert them to PNG or JPG."
        )

    if not images:
        raise ValueError(
            f"No valid PNG/JPG images found in background directory: {background_dir}"
        )

    return images


def select_background_pool(
    bg_images: List[Path],
    background_mode: BackgroundMode,
    background_image_number: int | None = None,
    background_cycle_start: int | None = None,
    background_cycle_end: int | None = None,
) -> List[Path]:
    if len(bg_images) == 1:
        return bg_images

    if background_mode == "fixed":
        assert background_image_number is not None
        if background_image_number > len(bg_images):
            raise ValueError(
                f"Requested background image {background_image_number}, "
                f"but only {len(bg_images)} images were found."
            )
        return [bg_images[background_image_number - 1]]

    if background_mode == "cycle":
        if background_cycle_start is None and background_cycle_end is None:
            return bg_images

        assert background_cycle_start is not None
        assert background_cycle_end is not None

        if background_cycle_end > len(bg_images):
            raise ValueError(
                f"Requested background cycle end {background_cycle_end}, "
                f"but only {len(bg_images)} images were found."
            )

        selected = bg_images[background_cycle_start - 1 : background_cycle_end]
        if not selected:
            raise ValueError("Background cycle range did not select any images.")
        return selected

    raise ValueError(f"Unsupported background mode: {background_mode}")


def resolve_background_image(
    bg_pool: List[Path],
    generated_slide_index: int,
) -> Path:
    if not bg_pool:
        raise ValueError("Background image pool is empty.")
    if len(bg_pool) == 1:
        return bg_pool[0]
    return bg_pool[generated_slide_index % len(bg_pool)]


def read_vocab_csv(csv_path: Path, min_rows: int = 10) -> List[VocabPair]:
    rows: List[VocabPair] = []
    with csv_path.open("r", encoding="utf-8-sig", newline="") as f:
        reader = csv.DictReader(f)
        if reader.fieldnames is None:
            raise ValueError("CSV file is empty or missing a header row.")

        field_map = {
            name.strip().lower(): name for name in reader.fieldnames if name is not None
        }
        if "term" not in field_map or "definition" not in field_map:
            raise ValueError(
                "CSV must contain headers for both Term and Definition "
                "(case-insensitive)."
            )

        term_key = field_map["term"]
        definition_key = field_map["definition"]

        for row in reader:
            term = (row.get(term_key) or "").strip()
            definition = (row.get(definition_key) or "").strip()
            if not term and not definition:
                continue
            if not term or not definition:
                raise ValueError(
                    f"Found incomplete row: Term={term!r}, Definition={definition!r}"
                )
            rows.append((term, definition))

    unique: List[VocabPair] = []
    seen = set()
    for pair in rows:
        if pair not in seen:
            unique.append(pair)
            seen.add(pair)

    if len(unique) < min_rows:
        raise ValueError(
            f"CSV must contain at least {min_rows} unique Term/Definition pairs; "
            f"found {len(unique)}."
        )

    return unique


def random_sets(
    vocab: List[VocabPair],
    set_count: int = 6,
    set_size: int = 10,
    seed: int = 42,
):
    if len(vocab) < set_size:
        raise ValueError(f"Need at least {set_size} unique items; got {len(vocab)}.")
    rng = random.Random(seed)
    return [rng.sample(vocab, set_size) for _ in range(set_count)]


def set_shape_fill_transparency(shape, transparency: float) -> None:
    if not 0.0 <= transparency <= 1.0:
        raise ValueError("transparency must be between 0.0 and 1.0")

    fill = shape.fill
    fill.solid()

    solid_fill = shape._element.spPr.solidFill
    if solid_fill is None:
        raise ValueError("Shape does not have a solid fill.")

    color_node = solid_fill.srgbClr
    if color_node is None:
        color_node = solid_fill.schemeClr
    if color_node is None:
        raise ValueError("Solid fill color node not found.")

    for child in list(color_node):
        if child.tag.endswith("alpha"):
            color_node.remove(child)

    alpha = OxmlElement("a:alpha")
    alpha.set("val", str(int((1.0 - transparency) * 100000)))
    color_node.append(alpha)


def apply_font(run, font_spec: Dict[str, Any], color: RGBColor) -> None:
    run.font.name = font_spec["name"]
    run.font.size = Pt(font_spec["size"])
    run.font.bold = font_spec["bold"]
    run.font.color.rgb = color


def add_default_background(slide, colors: Dict[str, RGBColor]):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colors["bg"]

    for left, top, width, height, color in [
        (Inches(-0.8), Inches(-0.6), Inches(2.6), Inches(2.2), colors["blue"]),
        (Inches(11.0), Inches(5.8), Inches(2.2), Inches(1.8), colors["coral"]),
    ]:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        set_shape_fill_transparency(shape, 0.82)
        shape.line.fill.background()


def add_image_background(prs, slide, image_path: Path):
    try:
        pic = slide.shapes.add_picture(
            str(image_path),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
    except Exception as exc:
        raise RuntimeError(f"Failed to load background image: {image_path}") from exc

    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)


def add_soft_overlay(prs, slide, transparency: float, colors: Dict[str, RGBColor]):
    overlay = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        prs.slide_width,
        prs.slide_height,
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = colors["white"]
    set_shape_fill_transparency(overlay, transparency)
    overlay.line.fill.background()


def apply_slide_background(
    prs,
    slide,
    colors: Dict[str, RGBColor],
    bg_image: Path | None = None,
):
    if bg_image is not None:
        add_image_background(prs, slide, bg_image)
    else:
        add_default_background(slide, colors)


def add_title_card(slide, colors: Dict[str, RGBColor]):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.7),
        Inches(0.9),
        Inches(11.9),
        Inches(4.9),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = colors["white"]
    card.line.color.rgb = colors["blue"]
    card.line.transparency = 0.85
    return card


def add_title_slide(
    prs,
    set_label: str,
    section_title: str,
    source_name: str,
    labels: Dict[str, str],
    style: Dict[str, Any],
    bg_image: Path | None = None,
):
    colors = style["colors"]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_background(prs, slide, colors, bg_image)

    if bg_image is not None:
        add_soft_overlay(
            prs,
            slide,
            transparency=style["overlay_transparency"],
            colors=colors,
        )

    text_left = Inches(0.9)
    text_top = Inches(1.45)
    text_width = Inches(11.0)
    text_height = Inches(3.9)

    if style["show_card"]:
        card = add_title_card(slide, colors)
        set_shape_fill_transparency(card, style["card_transparency"])
        text_left = Inches(1.1)
        text_top = Inches(1.55)
        text_width = Inches(10.6)
        text_height = Inches(3.5)

    pill = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.9),
        Inches(0.7),
        Inches(1.8),
        Inches(0.45),
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = colors["white"]
    set_shape_fill_transparency(pill, 0.15)
    pill.line.color.rgb = colors["blue"]

    tf = pill.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = set_label.upper()
    apply_font(run, style["pill_font"], colors["blue"])

    box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = labels["game_title"]
    apply_font(r, style["main_font"], colors["text"])

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = section_title
    apply_font(r2, style["section_font"], colors["blue"])

    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.LEFT
    r3 = p3.add_run()
    r3.text = f"{labels['source_prefix']} {source_name}"
    apply_font(r3, style["body_font"], colors["subtext"])


def fit_font_size(text: str) -> int:
    n = len(text)
    if n <= 10:
        return 34
    if n <= 18:
        return 30
    if n <= 26:
        return 26
    if n <= 36:
        return 22
    return 20


def resolve_vocab_primary_font(base_font: Dict[str, Any], text: str) -> Dict[str, Any]:
    resolved = dict(base_font)
    resolved["size"] = fit_font_size(text)
    return resolved


def add_vocab_slide(
    prs,
    primary_text: str,
    secondary_text: str | None,
    style: Dict[str, Any],
    bg_image: Path | None = None,
):
    colors = style["colors"]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_background(prs, slide, colors, bg_image)

    if bg_image is not None:
        add_soft_overlay(
            prs,
            slide,
            transparency=style["overlay_transparency"],
            colors=colors,
        )

    text_left = Inches(1.1)
    text_top = Inches(1.55)
    text_width = Inches(11.1)
    text_height = Inches(4.3)

    if style["show_card"]:
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.7),
            Inches(1.0),
            Inches(11.9),
            Inches(5.4),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = colors["white"]
        set_shape_fill_transparency(card, style["card_transparency"])
        card.line.color.rgb = colors["blue"]
        card.line.transparency = 0.85
    else:
        text_left = Inches(0.9)
        text_top = Inches(1.2)
        text_width = Inches(11.5)
        text_height = Inches(4.9)

    box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.clear()

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = primary_text
    apply_font(
        r1,
        resolve_vocab_primary_font(style["primary_font"], primary_text),
        colors["text"],
    )

    if secondary_text:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(12)
        r2 = p2.add_run()
        r2.text = secondary_text
        apply_font(r2, style["secondary_font"], colors["subtext"])


def build_pptx(
    csv_path: Path,
    output_path: Path,
    style_config: Dict[str, Any],
    set_count: int = 6,
    set_size: int = 10,
    seed: int = 42,
    primary_side: PrimarySide = "term",
    show_alternate: bool = True,
    export_selected_terms: bool = False,
    background_dir: str | Path | None = None,
    background_mode: BackgroundMode = "cycle",
    background_image_number: int | None = None,
    background_cycle_start: int | None = None,
    background_cycle_end: int | None = None,
    title_slide_overlay_transparency: float | None = None,
    vocab_slide_overlay_transparency: float | None = None,
    show_title_card: bool | None = None,
    title_card_transparency: float | None = None,
    show_vocab_card: bool | None = None,
    vocab_card_transparency: float | None = None,
):
    vocab = read_vocab_csv(csv_path, min_rows=set_size)
    sets = random_sets(vocab, set_count=set_count, set_size=set_size, seed=seed)

    bg_pool: List[Path] = []
    if background_dir is not None:
        bg_dir = resolve_background_dir(background_dir)
        bg_images = list_background_images(bg_dir)
        bg_pool = select_background_pool(
            bg_images,
            background_mode=background_mode,
            background_image_number=background_image_number,
            background_cycle_start=background_cycle_start,
            background_cycle_end=background_cycle_end,
        )

    labels = style_config["labels"]
    title_style = resolve_title_slide_style(style_config)
    vocab_style = resolve_vocab_slide_style(style_config)

    if title_slide_overlay_transparency is not None:
        title_style["overlay_transparency"] = title_slide_overlay_transparency
    if vocab_slide_overlay_transparency is not None:
        vocab_style["overlay_transparency"] = vocab_slide_overlay_transparency
    if show_title_card is not None:
        title_style["show_card"] = show_title_card
    if title_card_transparency is not None:
        title_style["card_transparency"] = title_card_transparency
    if show_vocab_card is not None:
        vocab_style["show_card"] = show_vocab_card
    if vocab_card_transparency is not None:
        vocab_style["card_transparency"] = vocab_card_transparency

    validate_visual_options(
        title_slide_overlay_transparency=title_style["overlay_transparency"],
        vocab_slide_overlay_transparency=vocab_style["overlay_transparency"],
        title_card_transparency=title_style["card_transparency"],
        vocab_card_transparency=vocab_style["card_transparency"],
    )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    source_name = csv_path.stem.replace("_", " ").replace("-", " ").title()
    rows = []
    generated_slide_index = 0

    for set_number, terms in enumerate(sets, start=1):
        title_bg = (
            resolve_background_image(bg_pool, generated_slide_index)
            if bg_pool
            else None
        )
        add_title_slide(
            prs,
            f"{labels['set_prefix']} {set_number}",
            f"{source_name} {labels['vocabulary_suffix']}",
            csv_path.name,
            labels=labels,
            style=title_style,
            bg_image=title_bg,
        )
        generated_slide_index += 1

        for term, definition in terms:
            if primary_side == "term":
                primary_text = term
                secondary_text = definition if show_alternate else None
            else:
                primary_text = definition
                secondary_text = term if show_alternate else None

            vocab_bg = (
                resolve_background_image(bg_pool, generated_slide_index)
                if bg_pool
                else None
            )
            add_vocab_slide(
                prs,
                primary_text,
                secondary_text,
                style=vocab_style,
                bg_image=vocab_bg,
            )
            rows.append((set_number, term, definition))
            generated_slide_index += 1

    prs.save(output_path)

    csv_out = None
    if export_selected_terms:
        csv_out = output_path.with_name(output_path.stem + "_selected_terms.csv")
        with csv_out.open("w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerow(["set_number", "term", "definition"])
            writer.writerows(rows)

    return output_path, csv_out


def generate_from_inputs(
    csv_file: str | Path,
    output: str | Path | None = None,
    style_config_file: str | Path | None = None,
    set_count: int = 6,
    set_size: int = 10,
    seed: int = 42,
    primary_side: PrimarySide = "term",
    show_alternate: bool = True,
    export_selected_terms: bool = False,
    background_dir: str | Path | None = None,
    background_mode: BackgroundMode = "cycle",
    background_image_number: int | None = None,
    background_cycle_start: int | None = None,
    background_cycle_end: int | None = None,
    title_slide_overlay_transparency: float | None = None,
    vocab_slide_overlay_transparency: float | None = None,
    show_title_card: bool | None = None,
    title_card_transparency: float | None = None,
    show_vocab_card: bool | None = None,
    vocab_card_transparency: float | None = None,
):
    validate_generation_options(set_count, set_size, primary_side)
    validate_background_options(
        background_dir=background_dir,
        background_mode=background_mode,
        background_image_number=background_image_number,
        background_cycle_start=background_cycle_start,
        background_cycle_end=background_cycle_end,
    )

    csv_path = resolve_csv_path(csv_file)
    output_path = resolve_output_path(csv_path, output)
    style_config = load_style_config(style_config_file)

    return build_pptx(
        csv_path,
        output_path,
        style_config=style_config,
        set_count=set_count,
        set_size=set_size,
        seed=seed,
        primary_side=primary_side,
        show_alternate=show_alternate,
        export_selected_terms=export_selected_terms,
        background_dir=background_dir,
        background_mode=background_mode,
        background_image_number=background_image_number,
        background_cycle_start=background_cycle_start,
        background_cycle_end=background_cycle_end,
        title_slide_overlay_transparency=title_slide_overlay_transparency,
        vocab_slide_overlay_transparency=vocab_slide_overlay_transparency,
        show_title_card=show_title_card,
        title_card_transparency=title_card_transparency,
        show_vocab_card=show_vocab_card,
        vocab_card_transparency=vocab_card_transparency,
    )
