import csv
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from src.mainkata.services.generator import (build_pptx, fit_font_size,
                                             generate_from_inputs,
                                             is_valid_image,
                                             list_background_images,
                                             random_sets, read_vocab_csv,
                                             resolve_background_dir,
                                             resolve_background_image,
                                             resolve_csv_path,
                                             resolve_output_path,
                                             select_background_pool,
                                             validate_background_options,
                                             validate_generation_options,
                                             validate_visual_options)


def _write_csv(path: Path, rows, headers=("Term", "Definition")):
    with path.open("w", encoding="utf-8", newline="") as f:
        writer = csv.writer(f)
        writer.writerow(headers)
        writer.writerows(rows)


def _write_image(path: Path, color=(120, 160, 200)):
    path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", (32, 24), color).save(path)


def _slide_text(slide) -> str:
    return "\n".join(
        shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text
    )


def _count_pictures(slide) -> int:
    return sum(
        1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    )


@pytest.fixture
def vocab_rows():
    return [(f"Term {i}", f"Definition {i}") for i in range(1, 13)]


@pytest.fixture
def vocab_csv(tmp_path: Path, vocab_rows):
    path = tmp_path / "vocab.csv"
    _write_csv(path, vocab_rows)
    return path


@pytest.fixture
def background_dir(tmp_path: Path) -> Path:
    bg_dir = tmp_path / "backgrounds"
    _write_image(bg_dir / "1.png", (255, 0, 0))
    _write_image(bg_dir / "2.jpg", (0, 255, 0))
    _write_image(bg_dir / "3.jpeg", (0, 0, 255))
    return bg_dir


class TestResolvePaths:
    def test_resolve_csv_path_returns_resolved_path_for_existing_file(
        self, vocab_csv: Path
    ):
        result = resolve_csv_path(vocab_csv)
        assert result == vocab_csv.resolve()

    def test_resolve_csv_path_raises_for_missing_file(self, tmp_path: Path):
        missing = tmp_path / "missing.csv"
        with pytest.raises(FileNotFoundError, match="CSV file not found"):
            resolve_csv_path(missing)

    def test_resolve_output_path_defaults_next_to_csv(self, vocab_csv: Path):
        result = resolve_output_path(vocab_csv)
        assert result == vocab_csv.with_name("vocab_vocab_sets.pptx")
        assert result.parent.exists()

    def test_resolve_output_path_creates_parent_directories(
        self, vocab_csv: Path, tmp_path: Path
    ):
        output = tmp_path / "nested" / "deeper" / "deck.pptx"
        result = resolve_output_path(vocab_csv, output)
        assert result == output.resolve()
        assert result.parent.exists()


class TestBackgroundHelpers:
    def test_resolve_background_dir_returns_resolved_directory(
        self, background_dir: Path
    ):
        assert resolve_background_dir(background_dir) == background_dir.resolve()

    def test_resolve_background_dir_raises_for_missing_directory(self, tmp_path: Path):
        with pytest.raises(FileNotFoundError, match="Background directory not found"):
            resolve_background_dir(tmp_path / "missing")

    def test_resolve_background_dir_raises_for_non_directory(self, tmp_path: Path):
        file_path = tmp_path / "file.txt"
        file_path.write_text("x", encoding="utf-8")
        with pytest.raises(ValueError, match="Background path is not a directory"):
            resolve_background_dir(file_path)

    def test_is_valid_image_detects_valid_and_invalid_files(self, tmp_path: Path):
        good = tmp_path / "good.png"
        bad = tmp_path / "bad.png"
        _write_image(good)
        bad.write_text("not an image", encoding="utf-8")
        assert is_valid_image(good) is True
        assert is_valid_image(bad) is False

    def test_list_background_images_returns_sorted_supported_images(
        self, background_dir: Path
    ):
        images = list_background_images(background_dir)
        assert [p.name for p in images] == ["1.png", "2.jpg", "3.jpeg"]

    def test_list_background_images_raises_for_invalid_image_file(self, tmp_path: Path):
        bg_dir = tmp_path / "backgrounds"
        bg_dir.mkdir()
        _write_image(bg_dir / "good.png")
        (bg_dir / "broken.jpg").write_text("broken", encoding="utf-8")
        with pytest.raises(ValueError, match="not valid PNG/JPG images"):
            list_background_images(bg_dir)

    def test_list_background_images_raises_when_no_supported_images(
        self, tmp_path: Path
    ):
        bg_dir = tmp_path / "backgrounds"
        bg_dir.mkdir()
        (bg_dir / "note.txt").write_text("hello", encoding="utf-8")
        with pytest.raises(ValueError, match="No valid PNG/JPG images found"):
            list_background_images(bg_dir)

    def test_select_background_pool_fixed_mode_returns_single_selected_image(
        self, background_dir: Path
    ):
        bg_images = list_background_images(background_dir)
        selected = select_background_pool(
            bg_images,
            background_mode="fixed",
            background_image_number=2,
        )
        assert [p.name for p in selected] == ["2.jpg"]

    def test_select_background_pool_fixed_mode_rejects_out_of_range_image(
        self, background_dir: Path
    ):
        bg_images = list_background_images(background_dir)
        with pytest.raises(ValueError, match="Requested background image 5"):
            select_background_pool(
                bg_images,
                background_mode="fixed",
                background_image_number=5,
            )

    def test_select_background_pool_cycle_mode_returns_all_when_no_range(
        self, background_dir: Path
    ):
        bg_images = list_background_images(background_dir)
        selected = select_background_pool(bg_images, background_mode="cycle")
        assert selected == bg_images

    def test_select_background_pool_cycle_mode_returns_requested_range(
        self, background_dir: Path
    ):
        bg_images = list_background_images(background_dir)
        selected = select_background_pool(
            bg_images,
            background_mode="cycle",
            background_cycle_start=2,
            background_cycle_end=3,
        )
        assert [p.name for p in selected] == ["2.jpg", "3.jpeg"]

    def test_select_background_pool_cycle_mode_rejects_out_of_range_end(
        self, background_dir: Path
    ):
        bg_images = list_background_images(background_dir)
        with pytest.raises(ValueError, match="Requested background cycle end 4"):
            select_background_pool(
                bg_images,
                background_mode="cycle",
                background_cycle_start=1,
                background_cycle_end=4,
            )

    def test_resolve_background_image_returns_single_image_for_single_pool(
        self, tmp_path: Path
    ):
        image_path = tmp_path / "only.png"
        _write_image(image_path)
        assert resolve_background_image([image_path], 10) == image_path

    def test_resolve_background_image_cycles_by_slide_index(self, background_dir: Path):
        bg_images = list_background_images(background_dir)
        assert resolve_background_image(bg_images, 0).name == "1.png"
        assert resolve_background_image(bg_images, 1).name == "2.jpg"
        assert resolve_background_image(bg_images, 4).name == "2.jpg"

    def test_resolve_background_image_rejects_empty_pool(self):
        with pytest.raises(ValueError, match="Background image pool is empty"):
            resolve_background_image([], 0)


class TestValidateGenerationOptions:
    def test_accepts_valid_options(self):
        validate_generation_options(2, 10, "term")
        validate_generation_options(1, 1, "definition")

    @pytest.mark.parametrize(
        "set_count,set_size,primary_side,error",
        [
            (0, 10, "term", "--sets must be at least 1."),
            (1, 0, "term", "--set-size must be at least 1."),
            (1, 10, "other", "--primary-side must be either 'term' or 'definition'."),
        ],
    )
    def test_rejects_invalid_options(self, set_count, set_size, primary_side, error):
        with pytest.raises(ValueError, match=error):
            validate_generation_options(set_count, set_size, primary_side)


class TestValidateBackgroundOptions:
    def test_accepts_none_background_dir(self):
        validate_background_options(None, "cycle", None, None, None)

    def test_accepts_valid_fixed_mode(self):
        validate_background_options("/tmp", "fixed", 1, None, None)

    def test_accepts_valid_cycle_mode_without_range(self):
        validate_background_options("/tmp", "cycle", None, None, None)

    def test_accepts_valid_cycle_mode_with_range(self):
        validate_background_options("/tmp", "cycle", None, 1, 3)

    @pytest.mark.parametrize(
        "background_mode,background_image_number,start,end,error",
        [
            ("bad", None, None, None, "--background-mode must be 'fixed' or 'cycle'."),
            (
                "fixed",
                None,
                None,
                None,
                "--background-image-number must be >= 1 when --background-mode=fixed.",
            ),
            ("fixed", 1, 1, 2, "cannot be used with --background-mode=fixed"),
            ("cycle", 1, None, None, "cannot be used with --background-mode=cycle"),
            ("cycle", None, 1, None, "must be provided together"),
            ("cycle", None, 0, 2, "--background-cycle-start must be at least 1."),
            ("cycle", None, 1, 0, "--background-cycle-end must be at least 1."),
            ("cycle", None, 3, 2, "cannot be greater than --background-cycle-end"),
        ],
    )
    def test_rejects_invalid_background_options(
        self, background_mode, background_image_number, start, end, error
    ):
        with pytest.raises(ValueError, match=error):
            validate_background_options(
                background_dir="/tmp/backgrounds",
                background_mode=background_mode,
                background_image_number=background_image_number,
                background_cycle_start=start,
                background_cycle_end=end,
            )


class TestValidateVisualOptions:
    def test_accepts_valid_visual_options(self):
        validate_visual_options(0.22, 0.33, 0.18, 0.25)

    @pytest.mark.parametrize(
        "values,error",
        [
            (
                (-0.1, 0.2, 0.2, 0.2),
                "--title-slide-overlay-transparency must be between 0.0 and 1.0.",
            ),
            (
                (0.2, 1.1, 0.2, 0.2),
                "--vocab-slide-overlay-transparency must be between 0.0 and 1.0.",
            ),
            (
                (0.2, 0.2, -0.1, 0.2),
                "--title-card-transparency must be between 0.0 and 1.0.",
            ),
            (
                (0.2, 0.2, 0.2, 1.1),
                "--vocab-card-transparency must be between 0.0 and 1.0.",
            ),
        ],
    )
    def test_rejects_invalid_visual_options(self, values, error):
        with pytest.raises(ValueError, match=error):
            validate_visual_options(*values)


class TestReadVocabCsv:
    def test_reads_case_insensitive_headers(self, tmp_path: Path):
        path = tmp_path / "mixed_headers.csv"
        _write_csv(
            path,
            [("hola", "hello"), ("adios", "goodbye")],
            headers=(" term ", "DEFINITION"),
        )
        result = read_vocab_csv(path, min_rows=2)
        assert result == [("hola", "hello"), ("adios", "goodbye")]

    def test_skips_fully_blank_rows(self, tmp_path: Path):
        path = tmp_path / "blank_rows.csv"
        _write_csv(path, [("a", "1"), ("", ""), ("b", "2")])
        result = read_vocab_csv(path, min_rows=2)
        assert result == [("a", "1"), ("b", "2")]

    def test_deduplicates_rows_preserving_first_occurrence(self, tmp_path: Path):
        path = tmp_path / "dupes.csv"
        _write_csv(path, [("a", "1"), ("b", "2"), ("a", "1"), ("c", "3")])
        result = read_vocab_csv(path, min_rows=3)
        assert result == [("a", "1"), ("b", "2"), ("c", "3")]

    def test_raises_for_missing_header_row(self, tmp_path: Path):
        path = tmp_path / "empty.csv"
        path.write_text("", encoding="utf-8")
        with pytest.raises(
            ValueError, match="CSV file is empty or missing a header row"
        ):
            read_vocab_csv(path)

    def test_raises_for_missing_required_headers(self, tmp_path: Path):
        path = tmp_path / "bad_headers.csv"
        _write_csv(path, [("a", "1")], headers=("Word", "Meaning"))
        with pytest.raises(
            ValueError, match="CSV must contain headers for both Term and Definition"
        ):
            read_vocab_csv(path, min_rows=1)

    def test_raises_for_incomplete_row(self, tmp_path: Path):
        path = tmp_path / "incomplete.csv"
        _write_csv(path, [("a", "1"), ("b", "")])
        with pytest.raises(ValueError, match="Found incomplete row"):
            read_vocab_csv(path, min_rows=1)

    def test_raises_when_unique_row_count_below_minimum(self, tmp_path: Path):
        path = tmp_path / "too_small.csv"
        _write_csv(path, [("a", "1"), ("a", "1"), ("b", "2")])
        with pytest.raises(
            ValueError, match="at least 3 unique Term/Definition pairs; found 2"
        ):
            read_vocab_csv(path, min_rows=3)


class TestRandomSets:
    def test_returns_reproducible_sets_for_same_seed(self, vocab_rows):
        result1 = random_sets(vocab_rows, set_count=3, set_size=4, seed=99)
        result2 = random_sets(vocab_rows, set_count=3, set_size=4, seed=99)
        assert result1 == result2

    def test_returns_requested_number_of_sets_and_items(self, vocab_rows):
        result = random_sets(vocab_rows, set_count=4, set_size=5, seed=7)
        assert len(result) == 4
        assert all(len(s) == 5 for s in result)
        assert all(len(set(s)) == 5 for s in result)

    def test_raises_when_vocab_smaller_than_set_size(self):
        with pytest.raises(ValueError, match="Need at least 5 unique items; got 2"):
            random_sets([("a", "1"), ("b", "2")], set_count=1, set_size=5)


class TestFitFontSize:
    @pytest.mark.parametrize(
        "text,expected",
        [
            ("short", 34),
            ("12345678901", 30),
            ("x" * 19, 26),
            ("x" * 27, 22),
            ("x" * 37, 20),
        ],
    )
    def test_returns_expected_size_by_length_band(self, text, expected):
        assert fit_font_size(text) == expected


class TestBuildPptx:
    def test_build_pptx_creates_expected_slide_count_and_file(
        self, vocab_csv: Path, tmp_path: Path
    ):
        output = tmp_path / "deck.pptx"
        pptx_path, csv_out = build_pptx(
            vocab_csv,
            output,
            set_count=2,
            set_size=3,
            seed=1,
            primary_side="term",
            show_alternate=True,
            export_selected_terms=False,
        )
        assert pptx_path == output
        assert pptx_path.exists()
        assert csv_out is None

        prs = Presentation(pptx_path)
        assert len(prs.slides) == 2 * (1 + 3)

    def test_build_pptx_exports_selected_terms_csv(
        self, vocab_csv: Path, tmp_path: Path
    ):
        output = tmp_path / "deck.pptx"
        pptx_path, csv_out = build_pptx(
            vocab_csv,
            output,
            set_count=2,
            set_size=3,
            seed=2,
            export_selected_terms=True,
        )
        assert pptx_path.exists()
        assert csv_out is not None
        assert csv_out.exists()

        with csv_out.open("r", encoding="utf-8", newline="") as f:
            rows = list(csv.reader(f))

        assert rows[0] == ["set_number", "term", "definition"]
        assert len(rows) == 1 + (2 * 3)
        assert set(row[0] for row in rows[1:]) == {"1", "2"}

    def test_build_pptx_uses_definition_as_primary_when_requested(
        self, vocab_csv: Path, tmp_path: Path
    ):
        output = tmp_path / "deck_definition_first.pptx"
        build_pptx(
            vocab_csv,
            output,
            set_count=1,
            set_size=1,
            seed=0,
            primary_side="definition",
            show_alternate=False,
        )

        prs = Presentation(output)
        vocab_slide = prs.slides[1]
        text_content = _slide_text(vocab_slide)
        assert "Definition" in text_content
        assert "Term" not in text_content

    def test_build_pptx_hides_alternate_text_when_disabled(
        self, vocab_csv: Path, tmp_path: Path
    ):
        output = tmp_path / "deck_hide_alt.pptx"
        build_pptx(
            vocab_csv, output, set_count=1, set_size=1, seed=0, show_alternate=False
        )
        prs = Presentation(output)
        vocab_slide = prs.slides[1]
        text_content = _slide_text(vocab_slide)
        assert "Term" in text_content
        assert text_content.count("Definition") == 0

    def test_build_pptx_includes_title_slide_text(
        self, vocab_csv: Path, tmp_path: Path
    ):
        output = tmp_path / "deck_title.pptx"
        build_pptx(vocab_csv, output, set_count=1, set_size=1, seed=0)
        prs = Presentation(output)
        title_text = _slide_text(prs.slides[0])
        assert "Vocabulary Games" in title_text
        assert "SET 1" in title_text
        assert "Source: vocab.csv" in title_text

    def test_build_pptx_can_hide_title_card(self, vocab_csv: Path, tmp_path: Path):
        output = tmp_path / "deck_no_title_card.pptx"
        build_pptx(
            vocab_csv,
            output,
            set_count=1,
            set_size=1,
            seed=0,
            show_title_card=False,
        )
        prs = Presentation(output)
        title_slide = prs.slides[0]
        assert "Vocabulary Games" in _slide_text(title_slide)

    def test_build_pptx_can_hide_vocab_card(self, vocab_csv: Path, tmp_path: Path):
        output = tmp_path / "deck_no_vocab_card.pptx"
        build_pptx(
            vocab_csv,
            output,
            set_count=1,
            set_size=1,
            seed=0,
            show_vocab_card=False,
        )
        prs = Presentation(output)
        vocab_slide = prs.slides[1]
        assert "Term" in _slide_text(vocab_slide)

    def test_build_pptx_applies_background_images_and_overlays(
        self, vocab_csv: Path, tmp_path: Path, background_dir: Path
    ):
        output = tmp_path / "deck_backgrounds.pptx"
        build_pptx(
            vocab_csv,
            output,
            set_count=1,
            set_size=2,
            seed=0,
            background_dir=background_dir,
            background_mode="cycle",
            title_slide_overlay_transparency=0.31,
            vocab_slide_overlay_transparency=0.47,
        )
        prs = Presentation(output)
        assert len(prs.slides) == 3
        assert _count_pictures(prs.slides[0]) == 1
        assert _count_pictures(prs.slides[1]) == 1
        assert _count_pictures(prs.slides[2]) == 1

    def test_build_pptx_uses_fixed_background_image_for_all_slides(
        self, vocab_csv: Path, tmp_path: Path, background_dir: Path
    ):
        output = tmp_path / "deck_fixed_bg.pptx"
        build_pptx(
            vocab_csv,
            output,
            set_count=1,
            set_size=2,
            seed=0,
            background_dir=background_dir,
            background_mode="fixed",
            background_image_number=2,
        )
        prs = Presentation(output)
        assert all(_count_pictures(slide) == 1 for slide in prs.slides)


class TestGenerateFromInputs:
    def test_generate_from_inputs_validates_and_creates_outputs(
        self, vocab_csv: Path, tmp_path: Path
    ):
        output = tmp_path / "generated" / "game_deck.pptx"
        pptx_path, csv_out = generate_from_inputs(
            vocab_csv,
            output=output,
            set_count=1,
            set_size=2,
            seed=5,
            primary_side="term",
            show_alternate=False,
            export_selected_terms=True,
            show_title_card=False,
            title_card_transparency=0.2,
            title_slide_overlay_transparency=0.3,
            show_vocab_card=True,
            vocab_card_transparency=0.4,
            vocab_slide_overlay_transparency=0.5,
        )
        assert pptx_path == output.resolve()
        assert pptx_path.exists()
        assert csv_out is not None and csv_out.exists()

    def test_generate_from_inputs_supports_background_options(
        self, vocab_csv: Path, tmp_path: Path, background_dir: Path
    ):
        output = tmp_path / "generated_with_bg.pptx"
        pptx_path, _ = generate_from_inputs(
            vocab_csv,
            output=output,
            set_count=1,
            set_size=2,
            background_dir=background_dir,
            background_mode="cycle",
            background_cycle_start=1,
            background_cycle_end=2,
        )
        assert pptx_path.exists()
        prs = Presentation(pptx_path)
        assert all(_count_pictures(slide) == 1 for slide in prs.slides)

    def test_generate_from_inputs_raises_before_path_resolution_for_bad_options(
        self, vocab_csv: Path
    ):
        with pytest.raises(ValueError, match="--sets must be at least 1"):
            generate_from_inputs(vocab_csv, set_count=0)

    def test_generate_from_inputs_raises_for_invalid_visual_options(
        self, vocab_csv: Path
    ):
        with pytest.raises(
            ValueError,
            match="--title-slide-overlay-transparency must be between 0.0 and 1.0.",
        ):
            generate_from_inputs(vocab_csv, title_slide_overlay_transparency=1.5)
